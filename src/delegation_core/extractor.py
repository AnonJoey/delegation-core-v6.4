"""
extractor.py — Convert file formats to plain text for vault ingestion.

Supported: .md/.markdown  .txt/.text  .csv  .html/.htm  .pdf  .docx  .xlsx  .pptx
Images are not supported — convert to .txt before dropping in the inbox.

All extractors return a string. Callers treat None or empty string as a failure
and route the file to _failed/.

Recursive split helpers (v0.3):
  extract_pages(path)    → list[str]  per-page text for PDFs; [] otherwise
"""

import csv as _csv
import logging
from pathlib import Path

logger = logging.getLogger("extractor")

SUPPORTED: frozenset[str] = frozenset({
    ".md", ".markdown",
    ".txt", ".text",
    ".csv",
    ".html", ".htm",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
})


def extract(path: Path) -> str | None:
    """
    Return extracted plain text from path.
    Returns None if the format is unsupported or extraction fails.
    """
    suffix = path.suffix.lower()
    _map = {
        ".md":   _text,
        ".markdown": _text,
        ".txt":  _text,
        ".text": _text,
        ".csv":  _csv_to_md,
        ".html": _html,
        ".htm":  _html,
        ".pdf":  _pdf,
        ".docx": _docx,
        ".xlsx": _xlsx,
        ".pptx": _pptx,
    }
    fn = _map.get(suffix)
    if fn is None:
        return None
    try:
        result = fn(path)
        return result if result and result.strip() else None
    except Exception as e:
        logger.warning("Extraction failed for %s: %s", path.name, e)
        return None


def format_label(path: Path) -> str:
    """Human-readable format name for a file path."""
    labels = {
        ".md": "Markdown", ".markdown": "Markdown",
        ".txt": "Text", ".text": "Text", ".csv": "CSV",
        ".html": "HTML", ".htm": "HTML", ".pdf": "PDF",
        ".docx": "Word", ".xlsx": "Excel", ".pptx": "PowerPoint",
    }
    return labels.get(path.suffix.lower(), path.suffix.upper().lstrip("."))


# ── extractors ────────────────────────────────────────────────────────────────

def _text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _csv_to_md(path: Path) -> str:
    rows = []
    with open(path, encoding="utf-8", errors="replace", newline="") as fh:
        reader = _csv.reader(fh)
        for i, row in enumerate(reader):
            if not any(c.strip() for c in row):
                continue
            cells = [c.replace("|", "\\|") for c in row]
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= 500:
                rows.append("| *(table truncated at 500 rows)* |")
                break
    return "\n".join(rows)


def _html(path: Path) -> str:
    from bs4 import BeautifulSoup
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "head", "nav", "footer", "aside"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def _pdf(path: Path) -> str:
    import pypdf
    reader = pypdf.PdfReader(str(path))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"[Page {i + 1}]\n{text.strip()}")
    if not pages:
        return (
            f"[Scanned PDF — no extractable text]\n"
            f"File: {path.name}\n"
            f"Pages: {len(reader.pages)}\n"
            "To make this searchable, export the text manually and drop it as a .txt file."
        )
    return "\n\n".join(pages)


def _docx(path: Path) -> str:
    import docx
    doc = docx.Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            cells = [str(c) if c is not None else "" for c in row]
            if not any(c.strip() for c in cells):
                continue
            rows.append("| " + " | ".join(cells) + " |")
            if i == 0:
                rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
            if i >= 200:
                rows.append("| *(sheet truncated at 200 rows)* |")
                break
        if rows:
            parts.append(f"## {sheet_name}\n\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


# ── Recursive split helpers (v0.3) ────────────────────────────────────────────

def extract_pages(path: Path) -> list[str]:
    """Return per-page text for PDFs. Returns [] for non-PDFs or on extraction error.

    Each element is the non-empty text of one page. Blank pages are dropped.
    Used by splitter.py to decide whether to split a PDF by page groups.
    """
    if path.suffix.lower() != ".pdf":
        return []
    try:
        import pypdf
        reader = pypdf.PdfReader(str(path))
        return [
            page.extract_text().strip()
            for page in reader.pages
            if (page.extract_text() or "").strip()
        ]
    except Exception as e:
        logger.warning("extract_pages failed for %s: %s", path.name, e)
        return []
